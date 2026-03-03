from fastapi import UploadFile, HTTPException
from typing import Dict, List
from io import BytesIO

import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document

def extract_pdf(file_bytes):
    text = ""
    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for page in doc:
        text += page.get_text("text")
        text += "\n\n---PAGE_BREAK---\n\n"

    return text

def extract_pptx(file_bytes):
    prs = Presentation(BytesIO(file_bytes))
    text = ""

    for i, slide in enumerate(prs.slides):
        text += f"\n\n---SLIDE_{i+1}---\n\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

def extract_docx(file_bytes):
    doc = Document(BytesIO(file_bytes))
    text = ""

    for para in doc.paragraphs:
        text += para.text + "\n"

    return text

def extract_txt(file_bytes):
    return file_bytes.decode("utf-8")

def clean_text(text: str):

    text = text.replace("\r", "")
    text = text.strip()
    return text

def split_into_sections(text: str, file_type: str) -> List[Dict]:
    sections = []

    if file_type == "pdf":
        raw_sections = text.split("---PAGE_BREAK---")
    elif file_type == "pptx":
        raw_sections = text.split("---SLIDE_")
    else:
        # For docx/txt fallback: split by double newline
        raw_sections = text.split("\n\n")

    for idx, section in enumerate(raw_sections):
        if not section.strip():
            continue

        sections.append({
            "section_id": idx + 1,
            "title": f"{file_type.upper()} Section {idx + 1}",
            "content": section.strip()
        })

    return sections

async def process_uploaded_file(file: UploadFile) -> Dict:
    """
    Accepts UploadFile
    Detects file type
    Extracts text
    Cleans text
    Splits into structured sections
    Returns structured JSON
    """

    filename = file.filename
    file_ext = filename.split(".")[-1].lower()

    file_bytes = await file.read()

    # -------- Route by type --------

    if file_ext == "pdf":
        raw_text = extract_pdf(file_bytes)

    elif file_ext == "pptx":
        raw_text = extract_pptx(file_bytes)

    elif file_ext == "docx":
        raw_text = extract_docx(file_bytes)

    elif file_ext in ["txt", "md"]:
        raw_text = extract_txt(file_bytes)

    else:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {file_ext}"
        )

    # -------- Clean --------
    cleaned_text = clean_text(raw_text)

    # -------- Structure --------
    sections = split_into_sections(cleaned_text, file_ext)

    # -------- Final JSON --------
    return {
        "filename": filename,
        "file_type": file_ext,
        "raw_text": cleaned_text,
        "sections": sections,
        "file_bytes": file_bytes,
    }
